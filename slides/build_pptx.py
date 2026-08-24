"""
build_pptx.py
Build the eight-minute talk as a PowerPoint deck on the Institute of Science
Tokyo 16:9 template (the layouts that carry the Science Tokyo logo).

  template : presentaion/format/パワーポイント_PowerPoint templates_16：9.pptx
             layout  1  A-タイトルスライド  (logo bottom-right, full-bleed art)
             layout 15  コンテンツ1-2       (logo top-right, free body area)
  design   : presentaion/slide_style_v3.md  — colour encodes argumentative role,
             the figure is the evidence and the slide text the interpretation,
             annotation budget 0–3 per plot, body ≤ 70 words, ≥ 16 pt.
  content  : docs/claim_freeze.md — seven slides, frozen order, frozen numbers.
  figures  : the `--style st` variants, built in the Science Tokyo palette:
                 python code/fig6_three_shifts.py --style st
                 python code/fig7_answer_talk.py  --style st
                 python code/fig4_tornado.py      --style st
                 python code/fig5_threshold.py    --style st

Run:
    python slides/build_pptx.py [path/to/template.pptx]
Out:
    slides/odmr_474nm_talk.pptx
"""

import pathlib
import sys

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

HERE = pathlib.Path(__file__).resolve().parent
REPO = HERE.parent
CODE = REPO / 'code'
TEMPLATE = pathlib.Path(sys.argv[1]) if len(sys.argv) > 1 else pathlib.Path(
    '/home/user/presentaion/format/パワーポイント_PowerPoint templates_16：9.pptx')
OUT = HERE / 'odmr_474nm_talk.pptx'

# slide_style_v3.md §3 — the palette is authoritative
SCIENCE_BLUE = RGBColor(0x1C, 0x31, 0x77)
DEEP_INDIGO = RGBColor(0x4B, 0x00, 0x82)
MIST_LAVENDER = RGBColor(0xE8, 0xEA, 0xF1)
MIDNIGHT_INK = RGBColor(0x10, 0x14, 0x26)
MUTED = RGBColor(0x6E, 0x74, 0x88)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

JP_FONT = 'Yu Gothic'          # the template's own face (see its slide 38)

TITLE_LAYOUT = 0               # A-タイトルスライド
BODY_LAYOUT = 14               # コンテンツ1-2

# content geometry, in inches
L = 0.55                       # left margin for body content
FIG_W = 7.70                   # figure column
COL_X = 8.45                   # right column
COL_W = 4.35
REF_Y = 6.92


# --------------------------------------------------------------------------
def _face(run, name=JP_FONT):
    """Set the latin AND east-asian typeface (python-pptx only sets latin)."""
    rPr = run._r.get_or_add_rPr()
    for tag in ('a:latin', 'a:ea', 'a:cs'):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {})
            rPr.append(el)
        el.set('typeface', name)


def _run(par, text, size, color, bold=False, sub=False, sup=False):
    r = par.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    if sub or sup:
        r._r.get_or_add_rPr().set('baseline', '-25000' if sub else '30000')
    _face(r)
    return r


def textbox(slide, x, y, w, h, blocks, size=18, color=MIDNIGHT_INK,
            space_after=10, line_spacing=1.25):
    """blocks: list of paragraphs; each is a str or a list of (text, **opts)."""
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, block in enumerate(blocks):
        par = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        par.line_spacing = line_spacing
        par.space_after = Pt(space_after)
        for piece in ([(block, {})] if isinstance(block, str) else block):
            text, opt = piece
            _run(par, text, opt.get('size', size), opt.get('color', color),
                 opt.get('bold', False), opt.get('sub', False),
                 opt.get('sup', False))
    return box


def set_ph(slide, idx, text, size=None, color=None, bold=None):
    ph = slide.placeholders[idx]
    tf = ph.text_frame
    tf.word_wrap = True
    par = tf.paragraphs[0]
    for r in list(par.runs):
        r._r.getparent().remove(r._r)
    r = par.add_run()
    r.text = text
    if size:
        r.font.size = Pt(size)
    if color:
        r.font.color.rgb = color
    if bold is not None:
        r.font.bold = bold
    _face(r)
    return ph


def drop_ph(slide, idx):
    ph = slide.placeholders[idx]
    ph._element.getparent().remove(ph._element)


def picture(slide, name, x, y, w):
    pic = slide.shapes.add_picture(str(CODE / name), Inches(x), Inches(y),
                                   width=Inches(w))
    return pic


def reference(slide, text):
    """§14: references bottom-left, 8–10 pt, muted."""
    textbox(slide, L, REF_Y, 9.5, 0.3, [text], size=9, color=MUTED,
            space_after=0)


def framed(slide, x, y, w, h, blocks, size=20, align_center=True):
    """§7: a thin frame fixes a claim (an equation, a condition) — not a card."""
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y),
                                Inches(w), Inches(h))
    sh.fill.background()
    sh.line.color.rgb = SCIENCE_BLUE
    sh.line.width = Pt(1.0)
    sh.shadow.inherit = False
    tf = sh.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = tf.margin_right = Inches(0.15)
    for i, block in enumerate(blocks):
        par = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if align_center:
            par.alignment = PP_ALIGN.CENTER
        for piece in ([(block, {})] if isinstance(block, str) else block):
            text, opt = piece
            _run(par, text, opt.get('size', size),
                 opt.get('color', MIDNIGHT_INK), opt.get('bold', False),
                 opt.get('sub', False), opt.get('sup', False))
    return sh


def table(slide, x, y, w, h, rows, col_w, header_fill=SCIENCE_BLUE,
          size=15, hi_row=None):
    """§12: minimal columns, navy header, no zebra, only the conclusion cell hot."""
    from pptx.enum.text import MSO_ANCHOR
    shape = slide.shapes.add_table(len(rows), len(rows[0]), Inches(x),
                                   Inches(y), Inches(w), Inches(h))
    tbl = shape.table
    tbl.first_row = True
    tbl.horz_banding = False
    for j, cw in enumerate(col_w):
        tbl.columns[j].width = Inches(cw)
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = tbl.cell(i, j)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.margin_left = cell.margin_right = Inches(0.08)
            cell.margin_top = cell.margin_bottom = Inches(0.04)
            tf = cell.text_frame
            tf.word_wrap = True
            par = tf.paragraphs[0]
            head = (i == 0)
            _run(par, str(val), size,
                 WHITE if head else MIDNIGHT_INK,
                 bold=head or (hi_row is not None and i == hi_row))
            if head:
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_fill
            elif hi_row is not None and i == hi_row:
                cell.fill.solid()
                cell.fill.fore_color.rgb = MIST_LAVENDER
            else:
                cell.fill.background()
    return tbl


# --------------------------------------------------------------------------
prs = Presentation(str(TEMPLATE))
# start from the empty template: drop the 48 example slides AND their parts,
# otherwise the new slides collide with the old part names in the package
for sld in list(prs.slides._sldIdLst):
    prs.part.drop_rel(sld.rId)
    prs.slides._sldIdLst.remove(sld)
layouts = prs.slide_masters[0].slide_layouts


def content_slide(title):
    s = prs.slides.add_slide(layouts[BODY_LAYOUT])
    set_ph(s, 0, title, size=26, color=SCIENCE_BLUE, bold=True)
    drop_ph(s, 10)
    return s


# ---- 1. title ------------------------------------------------------------
s1 = prs.slides.add_slide(layouts[TITLE_LAYOUT])
set_ph(s1, 0, '120 GPa の ODMR は 474 nm で測るべきである', size=32, bold=True)
set_ph(s1, 10, '緑 532 nm 比で感度 ×2.7、磁気マップ取得時間 ×7', size=17)
set_ph(s1, 11, 'R. Abe   Institute of Science Tokyo', size=15)
set_ph(s1, 12, '2026-08', size=12, color=SCIENCE_BLUE)
s1.notes_slide.notes_text_frame.text = (
    '0:00–0:40  結論を先に置く。市販 473 nm DPSS で最適の 0.03% 以内。')

# ---- 2. why sensitivity --------------------------------------------------
s2 = content_slide('感度が空間分解能と (P, T) 点数を決める')
framed(s2, L, 1.55, 6.6, 1.15,
       [[('測定時間  ∝  画素数 / η', {'size': 26}),
         ('B', {'size': 26, 'sub': True}), ('2', {'size': 26, 'sup': True})]])
textbox(s2, L, 3.15, 7.1,ac := 3.2,
        [[('水素化物超伝導の決定実験は', {}),
          ('局所磁気マップ', {'bold': True}),
          ('になった。', {})],
         [('CeH', {}), ('9', {'sub': True}),
          (' の Meissner 効果は、超伝導領域がミクロン'
           'スケールで強く不均一であることを示した。', {})],
         [('1 回のロードに数週間、1 回のアンビル破損で全損する領域では、', {}),
          ('感度は空間分解能・視野・到達できる (P, T) 点数を買う通貨', {'bold': True}),
          ('である。', {})]], size=18)
textbox(s2, COL_X + 0.6, 1.55, COL_W - 0.2, 1.2,
        [[('×2.7', {'size': 40, 'bold': True, 'color': SCIENCE_BLUE})],
         [('感度', {'size': 15, 'color': MUTED})]], space_after=2)
textbox(s2, COL_X + 0.6, 2.95, COL_W - 0.2, 1.2,
        [[('×7', {'size': 40, 'bold': True, 'color': SCIENCE_BLUE})],
         [('マップ 1 枚の取得時間', {'size': 15, 'color': MUTED})]], space_after=2)
reference(s2, 'P. Bhattacharyya et al., Nature 627, 73 (2024);  '
              'J. F. Barry et al., Rev. Mod. Phys. 92, 015004 (2020)')
s2.notes_slide.notes_text_frame.text = (
    '0:40–1:40  マップ時間 ∝ 画素数/η²。感度は指標ではなく通貨。')

# ---- 3. the two edges ----------------------------------------------------
s3 = content_slide('120 GPa では吸収帯とイオン化端が同時に青へ動く')
picture(s3, 'talk_three_shifts_st.png', L - 0.10, 1.45, FIG_W)
textbox(s3, COL_X, 1.55, COL_W, 4.6,
        [[('ZPL が +0.40 eV 青方偏移し、Huang–Rhys 因子が 3.08 → 4.61 に増える。'
           '吸収極大は 586 → 474 nm。', {})],
         [('同時に IP(³A₂) が 2.68 → 3.06 eV に上がる。3.06 eV = 405 nm。'
           'これより青は NV⁻ を直接 NV⁰ に変える。', {})],
         [('120 GPa では σ(473) は σ(532) の ', {}),
          ('10 倍', {'bold': True, 'color': SCIENCE_BLUE}), ('。', {})],
         [('η ∝ Δν/(C√R) はコントラストに線形、光子レートに √ で効く。'
           '釣り合いの位置は量的問題である。', {})]], size=17)
reference(s3, 'K. O. Ho, C. Dailledouze et al., arXiv:2606.02399 (2026) — DFT + DAC。'
              '灰色の常圧曲線は単一有効フォノン近似（50 GPa 以下では粗い）。')
s3.notes_slide.notes_text_frame.text = (
    '1:40–3:00  2 つの端が窓を挟み込む。常圧曲線は近似が粗い旨を先に断る。')

# ---- 4. the answer -------------------------------------------------------
s4 = content_slide('最適は 474 nm、5% 許容窓は 462–486 nm')
picture(s4, 'talk_answer_120GPa_st.png', L - 0.10, 1.45, FIG_W)
table(s4, COL_X, 1.45, COL_W, 1.9,
      [['励起線', 'η/η（最適）', ''],
       ['405 nm', '3.67', 'イオン化端'],
       ['473 nm', '1.00', '市販 DPSS'],
       ['532 nm', '2.66', '従来の既定']],
      col_w=[1.35, 1.15, 1.85], hi_row=2)
textbox(s4, COL_X, 3.75, COL_W, 2.4,
        [[('473 nm DPSS で最適の 0.03% 以内。', {'bold': True})],
         [('λ', {}), ('opt', {'sub': True}),
          (' は ZPL に追随して 0.5 nm/GPa で動く。100 GPa では 484 nm。', {})],
         [('97–172 GPa を通す実験なら 488 nm と 473 nm の 2 本立てが有利。', {})]],
        size=17)
reference(s4, '帯は現象論定数のモンテカルロ 16–84%。'
              'λopt = 474.0 +3.7 / −4.3 nm。')
s4.notes_slide.notes_text_frame.text = (
    '3:00–4:20  答え。窓の広さと 473 nm の位置を見せる。')

# ---- 5. tornado ----------------------------------------------------------
s5 = content_slide('最適波長を動かすのは測定量だけである')
picture(s5, 'tornado_lambda_opt_120GPa_st.png', L - 0.10, 1.35, FIG_W)
textbox(s5, COL_X, 1.55, COL_W, 4.6,
        [[('励起状態イオン化と再結合はどちらも σ', {}), ('abs', {'sub': True}),
          (' に比例するので、定常 NV⁻ 分率 f₋ から 0.2% の精度で相殺する。', {})],
         [('固定パワーでの最適は ', {}),
          ('吸収極大そのもの', {'bold': True, 'color': SCIENCE_BLUE}), ('。', {})],
         [('校正でフィットする 9 個の定数は 1 つも λ', {}), ('opt', {'sub': True}),
          (' を動かさない。', {})],
         [('追加調整なしで、0–150 GPa・4 波長にわたる既発表 26 件を再現。', {})]],
        size=17)
reference(s5, '各入力をモンテカルロ範囲の端に振り、他を公称値に固定したときの λopt の変化。')
s5.notes_slide.notes_text_frame.text = (
    '4:20–5:20  相殺 → 現象論非依存。ゼロが並ぶ絵が主張そのもの。')

# ---- 6. threshold --------------------------------------------------------
s6 = content_slide('約 70 GPa で緑と青が入れ替わる')
picture(s6, 'threshold_green_blue_st.png', L - 0.10, 1.45, FIG_W)
textbox(s6, COL_X, 1.55, COL_W, 4.6,
        [[('50 GPa で青と緑を比較した既報は、', {}),
          ('青の明確な優位を見つけなかった', {'bold': True}), ('。', {})],
         [('モデルはそこで 0.59、つまり緑が有利と答える。'
           '見つからなかったことが再現になっている。', {})],
         [('単一圧力の比較に乗る波長依存の系統誤差は圧力に依らないので、'
           '比が 1 を横切る位置を作ることも動かすこともできない。', {})],
         [('これが本研究を反証する測定である。',
           {'bold': True, 'color': SCIENCE_BLUE})]], size=17)
reference(s6, 'P. Bhattacharyya, PhD thesis, UC Berkeley (2022), Sec. 6.3.')
s6.notes_slide.notes_text_frame.text = (
    '5:20–6:40  反証条件。符号変化は単一圧力の系統誤差では作れない。')

# ---- 7. scope ------------------------------------------------------------
s7 = content_slide('474 nm が成り立つ条件、成り立たない条件')
table(s7, L, 1.35, 12.2, 3.9,
      [['境界', '条件', '外れたとき'],
       ['応力', '準静水圧 (micropillar, α ≈ 0.95)',
        '平坦キュレット (α = 0.56) なら 508 nm。その形状は 40–50 GPa で '
        'ODMR コントラストを失う — micropillar は前提条件'],
       ['圧力', '約 70 GPa 以上', 'それ以下では緑が有利。推奨が反転する'],
       ['パワー', '低励起 (u ≲ 0.1)',
        '飽和が始まると最適は青へ動く。u = 0.3 で 473 nm 固定のコストは ×1.6'],
       ['アンビル', '低窒素アンビルへのイオン注入',
        'type Ib バルクは 500 nm 以下を強く吸収し、青を送れない']],
      col_w=[1.3, 3.5, 7.4], size=14)
textbox(s7, L, 5.55, 12.2, 1.1,
        [[('検出帯も ZPL とともに青へ動く。常圧で選んだパスバンドは 120 GPa で '
           '発光の 26% しか拾わない。青励起にして初めて検出カットオンも動かせ、'
           '×1.6 が上乗せされる — 合計 ×4。', {})],
         [('120 GPa では 532 nm を 473 nm DPSS に替える。'
           '同じ磁気マップが 1/7 の時間で撮れる。',
           {'bold': True, 'color': SCIENCE_BLUE, 'size': 19})]], size=16)
s7.notes_slide.notes_text_frame.text = (
    '6:40–8:00  適用条件とまとめ。校正は 474 nm の検証ではなく絶対感度・'
    'パワー依存・f₋ の直接測定のために行う。')

prs.save(str(OUT))
print(f'wrote {OUT}  ({OUT.stat().st_size / 1024:.0f} kB, '
      f'{len(prs.slides.__iter__.__self__._sldIdLst)} slides)')
